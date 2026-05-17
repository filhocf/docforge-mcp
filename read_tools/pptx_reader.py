"""Read and extract content from existing PPTX files."""

from pathlib import Path
from pptx import Presentation


def read_pptx(file_path: str) -> str:
    """Extract all text from a PPTX file."""
    prs = Presentation(file_path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def get_pptx_info(file_path: str) -> dict:
    """Get metadata and statistics from a PPTX file."""
    prs = Presentation(file_path)
    core = prs.core_properties

    slides_count = len(prs.slides)
    total_shapes = sum(len(slide.shapes) for slide in prs.slides)

    return {
        "file": Path(file_path).name,
        "title": core.title or "",
        "author": core.author or "",
        "created": str(core.created) if core.created else "",
        "modified": str(core.modified) if core.modified else "",
        "slides": slides_count,
        "total_shapes": total_shapes,
        "slide_width": str(prs.slide_width),
        "slide_height": str(prs.slide_height),
    }


def get_pptx_slides(file_path: str) -> list[dict]:
    """List all slides with index, layout, and text content."""
    prs = Presentation(file_path)
    result = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        result.append({
            "index": i,
            "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "shapes": len(slide.shapes),
            "text_content": texts,
        })
    return result

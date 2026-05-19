"""PowerPoint slide builder class.

This module provides the PowerpointPresentation class which builds slides
from structured data, combining helper mixins for text, tables, images, etc.
"""

import io
import logging
from typing import Any, Dict, List

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from template_utils import find_pptx_templates

from .chart_utils import ChartDataError, add_chart_to_slide
from .constants import (
    CONTENT_LAYOUT,
    DEFAULT_CAPTION_FONT_SIZE,
    DEFAULT_QUOTE_FONT_SIZE,
    DEFAULT_SUBTITLE_FONT_SIZE,
    SECTION_LAYOUT,
    TABLE_HEADER_FILL,
    TITLE_LAYOUT,
    TWO_COLUMN_LAYOUT,
    TWO_COLUMN_TEXT_LAYOUT,
)
from .helpers import (
    ImageHelperMixin,
    SlideHelperMixin,
    TableHelperMixin,
    TextHelperMixin,
    parse_color,
    parse_table_data,
)

logger = logging.getLogger(__name__)


def _load_templates():
    """Load presentation templates for 4:3 and 16:9 formats.

    Returns:
        Tuple of (path_4_3, path_16_9) template paths.
    """
    t43, t169 = find_pptx_templates()
    if not t43 or not t169:
        logger.info("One or more PPT templates missing; using PowerPoint defaults")
    return t43, t169


class PowerpointPresentation(SlideHelperMixin, TextHelperMixin, TableHelperMixin, ImageHelperMixin):
    """Builder class for creating PowerPoint presentations from structured data."""

    def __init__(self, slides: List[Dict[str, Any]], format: str):
        """Initialize and build presentation.

        Args:
            slides: List of slide dictionaries.
            format: Presentation format ("4:3" or "16:9").
        """
        logger.info(f"Initializing PowerPoint: slides={len(slides)}, format={format}")

        if not slides:
            raise ValueError("At least one slide is required")

        self.presentation = self._create_presentation(format)
        self._remove_default_slide()
        self._build_slides(slides)

    def _create_presentation(self, format: str) -> Presentation:
        """Create presentation with appropriate template.

        Args:
            format: "4:3" or "16:9".

        Returns:
            Presentation object.
        """
        template_4_3, template_16_9 = _load_templates()
        template = template_16_9 if format == "16:9" else template_4_3

        if template:
            try:
                return Presentation(template)
            except Exception as e:
                logger.error(f"Failed to load template: {e}")

        logger.warning(f"Using default PowerPoint template for {format}")
        return Presentation()

    def _remove_default_slide(self) -> None:
        """Remove default slide if present."""
        if len(self.presentation.slides) > 0:
            try:
                slide = self.presentation.slides[0]
                self.presentation.slides.element.remove(slide.element)
                logger.debug("Removed default slide")
            except Exception as e:
                logger.debug(f"Could not remove default slide: {e}")

    def _build_slides(self, slides: List[Dict[str, Any]]) -> None:
        """Build all slides from data.

        Args:
            slides: List of slide dictionaries.
        """
        slide_builders = {
            "title": self._build_title_slide,
            "section": self._build_section_slide,
            "content": self._build_content_slide,
            "table": self._build_table_slide,
            "image": self._build_image_slide,
            "two_column": self._build_two_column_slide,
            "chart": self._build_chart_slide,
            "quote": self._build_quote_slide,
        }

        logger.info(f"Building {len(slides)} slides")

        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get("slide_type", "")
            builder = slide_builders.get(slide_type)

            if builder:
                try:
                    logger.debug(f"Building slide {i}: type={slide_type}")
                    builder(slide_data)
                except Exception as e:
                    logger.error(f"Failed to create slide {i}: {e}")
                    raise ValueError(f"Error creating slide {i} ({slide_type}): {e}")
            else:
                logger.warning(f"Unknown slide type '{slide_type}' at index {i}")

    # -------------------------------------------------------------------------
    # Slide Builders
    # -------------------------------------------------------------------------

    def _build_title_slide(self, data: Dict[str, Any]) -> None:
        """Build a title slide with title and author."""
        layout = self.presentation.slide_layouts[TITLE_LAYOUT]
        slide = self.presentation.slides.add_slide(layout)

        if len(slide.placeholders) > 0:
            slide.placeholders[0].text = data.get("slide_title", "")
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = data.get("author", "")

        self._add_speaker_notes(slide, data.get("speaker_notes"))

    def _build_section_slide(self, data: Dict[str, Any]) -> None:
        """Build a section divider slide."""
        layout = self.presentation.slide_layouts[SECTION_LAYOUT]
        slide = self.presentation.slides.add_slide(layout)

        if len(slide.placeholders) > 0:
            slide.placeholders[0].text = data.get("slide_title", "")

        self._add_speaker_notes(slide, data.get("speaker_notes"))

    def _build_content_slide(self, data: Dict[str, Any]) -> None:
        """Build a content slide with bullet points."""
        layout = self.presentation.slide_layouts[CONTENT_LAYOUT]
        slide = self.presentation.slides.add_slide(layout)

        # Title
        if len(slide.placeholders) > 0:
            slide.placeholders[0].text = data.get("slide_title", "")

        # Bullet points
        slide_text = data.get("slide_text", [])
        if slide_text and len(slide.placeholders) > 1:
            placeholder = slide.placeholders[1]
            placeholder.text = ""

            for i, item in enumerate(slide_text):
                para = placeholder.text_frame.paragraphs[0] if i == 0 else placeholder.text_frame.add_paragraph()
                para.text = item.get("text", "")
                para.alignment = PP_ALIGN.LEFT
                para.level = max(0, int(item.get("indentation_level", 1)) - 1)

        self._add_speaker_notes(slide, data.get("speaker_notes"))

    def _build_table_slide(self, data: Dict[str, Any]) -> None:
        """Build a table slide with styled table using Title and Content layout."""
        title = data.get("slide_title", "")
        slide, left, top, width, height = self._add_title_content_slide(title)

        # Table
        table_data = parse_table_data(data.get("table_data", []))
        if not table_data:
            logger.warning("No table data provided")
            return

        header_color = parse_color(
            data.get("header_color", "4172C4"),
            TABLE_HEADER_FILL
        )

        self._create_styled_table(
            slide,
            table_data,
            left=left,
            top=top,
            width=width,
            height=height,
            header_color=header_color,
            alternate_rows=data.get("alternate_rows", True)
        )

        self._add_speaker_notes(slide, data.get("speaker_notes"))

    def _build_image_slide(self, data: Dict[str, Any]) -> None:
        """Build a slide with an image from URL using Title and Content layout."""
        title = data.get("slide_title", "")
        slide, left, top, width, height = self._add_title_content_slide(title)

        # Image
        image_url = data.get("image_url", "")
        caption = data.get("image_caption", "")

        max_height = height - (Inches(0.6) if caption else 0)

        if image_url:
            picture = self._add_image_from_url(
                slide, image_url,
                left=left,
                top=top,
                max_width=width,
                max_height=max_height
            )

            if picture and caption:
                self._add_text_box(
                    slide, caption,
                    left=left,
                    top=picture.top + picture.height + Inches(0.1),
                    width=width,
                    height=Inches(0.5),
                    font_size=DEFAULT_CAPTION_FONT_SIZE,
                    italic=True,
                    alignment=PP_ALIGN.CENTER
                )
            elif not picture:
                self._add_image_placeholder(
                    slide, "Image could not be loaded",
                    left, top + Inches(1), width
                )

        self._add_speaker_notes(slide, data.get("speaker_notes"))

    def _build_two_column_slide(self, data: Dict[str, Any]) -> None:
        """Build a slide with two text columns using built-in PowerPoint layouts.

        Uses TWO_COLUMN_TEXT_LAYOUT (Comparison) if subheaders are provided,
        otherwise uses TWO_COLUMN_LAYOUT (Two Content).

        Placeholder indices:
        - Two Content (3): idx 0=Title, 1=Left content, 2=Right content
        - Comparison (4): idx 0=Title, 1=Left subheader, 2=Left content, 3=Right subheader, 4=Right content
        """
        left_heading = data.get("left_heading", "")
        right_heading = data.get("right_heading", "")
        has_subheaders = bool(left_heading or right_heading)

        # Choose layout based on whether subheaders are needed
        if has_subheaders:
            layout = self.presentation.slide_layouts[TWO_COLUMN_TEXT_LAYOUT]
        else:
            layout = self.presentation.slide_layouts[TWO_COLUMN_LAYOUT]

        slide = self.presentation.slides.add_slide(layout)

        # Fill placeholders based on layout type
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx

            # Title placeholder (idx 0) - both layouts
            if idx == 0:
                title = data.get("slide_title", "")
                if title:
                    shape.text = title

            elif has_subheaders:
                # Comparison layout indices
                if idx == 1:  # Left subheader
                    if left_heading:
                        shape.text = left_heading
                elif idx == 2:  # Left content
                    self._fill_placeholder_with_bullets(shape, data.get("left_column", []))
                elif idx == 3:  # Right subheader
                    if right_heading:
                        shape.text = right_heading
                elif idx == 4:  # Right content
                    self._fill_placeholder_with_bullets(shape, data.get("right_column", []))
            else:
                # Two Content layout indices
                if idx == 1:  # Left content
                    self._fill_placeholder_with_bullets(shape, data.get("left_column", []))
                elif idx == 2:  # Right content
                    self._fill_placeholder_with_bullets(shape, data.get("right_column", []))

        self._add_speaker_notes(slide, data.get("speaker_notes"))


    def _build_chart_slide(self, data: Dict[str, Any]) -> None:
        """Build a slide with a chart using Title and Content layout."""
        title = data.get("slide_title", "")
        slide, left, top, width, height = self._add_title_content_slide(title)

        # Chart
        chart_data = data.get("chart_data", {})
        if not chart_data:
            self._add_text_box(
                slide, "[No chart data provided]",
                left, top, width, Inches(1),
                alignment=PP_ALIGN.CENTER
            )
            return

        try:
            add_chart_to_slide(
                slide,
                chart_type=data.get("chart_type", "bar"),
                chart_data=chart_data,
                left=left,
                top=top,
                width=width,
                height=height,
                has_legend=data.get("has_legend", True),
                legend_position=data.get("legend_position", "right")
            )
        except ChartDataError as e:
            logger.error(f"Chart error: {e}")
            self._add_text_box(
                slide, f"[Chart error: {e}]",
                left, top, width, Inches(1),
                alignment=PP_ALIGN.CENTER
            )

        self._add_speaker_notes(slide, data.get("speaker_notes"))


    def _build_quote_slide(self, data: Dict[str, Any]) -> None:
        """Build a quote/citation slide using Title and Content layout."""
        title = data.get("slide_title", "")
        slide, left, top, width, height = self._add_title_content_slide(title)

        # Quote
        quote_text = data.get("quote_text", "")
        quote_author = data.get("quote_author", "")

        quote_box = slide.shapes.add_textbox(left, top, width, height)
        tf = quote_box.text_frame
        tf.word_wrap = True

        # Quote text
        para = tf.paragraphs[0]
        para.text = f'"{quote_text}"'
        para.font.size = DEFAULT_QUOTE_FONT_SIZE
        para.font.italic = True
        para.alignment = PP_ALIGN.CENTER

        # Author
        if quote_author:
            author_para = tf.add_paragraph()
            author_para.text = f"— {quote_author}"
            author_para.font.size = DEFAULT_SUBTITLE_FONT_SIZE
            author_para.font.bold = True
            author_para.alignment = PP_ALIGN.CENTER
            author_para.space_before = Pt(24)

        self._add_speaker_notes(slide, data.get("speaker_notes"))

    # -------------------------------------------------------------------------
    # Output
    # -------------------------------------------------------------------------

    def save(self) -> io.BytesIO:
        """Save presentation to a BytesIO object.

        Returns:
            BytesIO containing the presentation.

        Raises:
            RuntimeError: If saving fails.
        """
        logger.info("Saving PowerPoint to memory buffer")
        try:
            buffer = io.BytesIO()
            self.presentation.save(buffer)
            buffer.seek(0)
            return buffer
        except Exception as e:
            logger.error("Failed to save PowerPoint presentation: %s", e, exc_info=True)
            raise RuntimeError(f"Failed to save presentation: {e}") from e


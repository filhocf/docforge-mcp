"""Advanced PPTX features implementation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree


def open_pptx_and_edit(file_path: str, slide_index: int, placeholder_map: dict[str, str], output_path: str | None = None) -> str:
    """Open existing PPTX and replace placeholder text in a slide.

    placeholder_map: {"{{title}}": "Real Title", "{{subtitle}}": "Real Sub"}
    """
    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})")

    slide = prs.slides[slide_index]
    replacements = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, value in placeholder_map.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)
                            replacements += 1

    save_path = output_path or file_path
    prs.save(save_path)
    return f"Replaced {replacements} placeholder(s) in slide {slide_index}. Saved to {save_path}"


def add_shape_to_slide(file_path: str, slide_index: int, shape_type: str = "rectangle",
                       left_inches: float = 1.0, top_inches: float = 1.0,
                       width_inches: float = 3.0, height_inches: float = 2.0,
                       text: str = "", fill_color: str | None = None,
                       output_path: str | None = None) -> str:
    """Add a shape to a slide. shape_type: rectangle, oval, rounded_rectangle, triangle."""
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    }
    if shape_type not in shape_map:
        raise ValueError(f"Unsupported shape: '{shape_type}'. Supported: {list(shape_map.keys())}")

    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]
    shape = slide.shapes.add_shape(
        shape_map[shape_type],
        Inches(left_inches), Inches(top_inches),
        Inches(width_inches), Inches(height_inches),
    )

    if text:
        shape.text = text

    if fill_color:
        color = fill_color.lstrip("#")
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(int(color[:2], 16), int(color[2:4], 16), int(color[4:6], 16))

    save_path = output_path or file_path
    prs.save(save_path)
    return f"Shape '{shape_type}' added to slide {slide_index}. Saved to {save_path}"


def duplicate_slide(file_path: str, slide_index: int, output_path: str | None = None) -> str:
    """Duplicate a slide (copy it to the end of the presentation)."""
    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy shapes
    for shape in source_slide.shapes:
        el = deepcopy(shape.element)
        new_slide.shapes._spTree.append(el)

    save_path = output_path or file_path
    prs.save(save_path)
    return f"Slide {slide_index} duplicated (now {len(prs.slides)} slides). Saved to {save_path}"


def apply_slide_template(file_path: str, slide_index: int, replacements: dict[str, str], output_path: str | None = None) -> str:
    """Apply template replacements to all text in a slide. Same as open_pptx_and_edit but clearer name."""
    return open_pptx_and_edit(file_path, slide_index, replacements, output_path)

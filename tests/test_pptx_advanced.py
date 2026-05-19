"""Tests for advanced PPTX features."""


import pytest
from pptx import Presentation

from pptx_tools.advanced import add_shape_to_slide, apply_slide_template, duplicate_slide, open_pptx_and_edit


@pytest.fixture
def template_pptx(tmp_path):
    path = str(tmp_path / "template.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "{{title}}"
    slide.placeholders[1].text = "{{subtitle}}"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Content"
    prs.save(path)
    return path


class TestPlaceholderReplace:
    def test_replace_placeholders(self, template_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        result = open_pptx_and_edit(template_pptx, 0, {"{{title}}": "Real Title", "{{subtitle}}": "Real Sub"}, out)
        assert "2" in result  # 2 replacements
        prs = Presentation(out)
        assert prs.slides[0].shapes.title.text == "Real Title"

    def test_invalid_slide_index(self, template_pptx):
        with pytest.raises(ValueError):
            open_pptx_and_edit(template_pptx, 99, {"x": "y"})


class TestAddShape:
    def test_add_rectangle(self, template_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        result = add_shape_to_slide(template_pptx, 0, "rectangle", text="Hello", output_path=out)
        assert "rectangle" in result
        prs = Presentation(out)
        shapes = prs.slides[0].shapes
        assert any(s.text == "Hello" for s in shapes if s.has_text_frame)

    def test_add_oval_with_color(self, template_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        result = add_shape_to_slide(template_pptx, 0, "oval", fill_color="#FF0000", output_path=out)
        assert "oval" in result

    def test_invalid_shape(self, template_pptx):
        with pytest.raises(ValueError, match="Unsupported"):
            add_shape_to_slide(template_pptx, 0, "hexagon")


class TestDuplicateSlide:
    def test_duplicate(self, template_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        result = duplicate_slide(template_pptx, 0, out)
        assert "3 slides" in result
        prs = Presentation(out)
        assert len(prs.slides) == 3

    def test_duplicate_invalid(self, template_pptx):
        with pytest.raises(ValueError):
            duplicate_slide(template_pptx, 99)


class TestApplyTemplate:
    def test_apply_template(self, template_pptx, tmp_path):
        out = str(tmp_path / "out.pptx")
        result = apply_slide_template(template_pptx, 0, {"{{title}}": "Applied"}, out)
        assert "1" in result
        prs = Presentation(out)
        assert prs.slides[0].shapes.title.text == "Applied"

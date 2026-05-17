"""Tests for read_tools — reading existing DOCX, XLSX, PPTX files."""

import os
import tempfile
import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from read_tools import (
    read_docx, get_docx_info, get_docx_paragraphs, get_docx_tables,
    read_xlsx, get_xlsx_info, get_xlsx_sheets,
    read_pptx, get_pptx_info, get_pptx_slides,
)


@pytest.fixture
def sample_docx(tmp_path):
    """Create a sample DOCX file for testing."""
    path = tmp_path / "test.docx"
    doc = Document()
    doc.core_properties.title = "Test Document"
    doc.core_properties.author = "Test Author"
    doc.add_heading("Chapter 1", level=1)
    doc.add_paragraph("This is the first paragraph.")
    doc.add_paragraph("This is the second paragraph.")
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(0, 2).text = "City"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "30"
    table.cell(1, 2).text = "NYC"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def sample_xlsx(tmp_path):
    """Create a sample XLSX file for testing."""
    path = tmp_path / "test.xlsx"
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Data"
    ws1.append(["Name", "Score", "Grade"])
    ws1.append(["Alice", 95, "A"])
    ws1.append(["Bob", 82, "B"])
    ws2 = wb.create_sheet("Summary")
    ws2.append(["Total", 177])
    wb.save(str(path))
    return str(path)


@pytest.fixture
def sample_pptx(tmp_path):
    """Create a sample PPTX file for testing."""
    path = tmp_path / "test.pptx"
    prs = Presentation()
    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Test Author"
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Welcome"
    slide1.placeholders[1].text = "Subtitle text"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Content Slide"
    prs.save(str(path))
    return str(path)


# === DOCX Tests ===

class TestDocxReader:
    def test_read_docx(self, sample_docx):
        text = read_docx(sample_docx)
        assert "Chapter 1" in text
        assert "first paragraph" in text
        assert "second paragraph" in text

    def test_get_docx_info(self, sample_docx):
        info = get_docx_info(sample_docx)
        assert info["title"] == "Test Document"
        assert info["author"] == "Test Author"
        assert info["paragraphs"] > 0
        assert info["tables"] == 1
        assert info["characters"] > 0

    def test_get_docx_paragraphs(self, sample_docx):
        paragraphs = get_docx_paragraphs(sample_docx)
        assert len(paragraphs) > 0
        assert paragraphs[0]["index"] == 0
        assert "style" in paragraphs[0]
        assert "text" in paragraphs[0]
        # Find heading
        heading = next(p for p in paragraphs if "Heading" in p["style"])
        assert heading["text"] == "Chapter 1"

    def test_get_docx_tables(self, sample_docx):
        tables = get_docx_tables(sample_docx)
        assert len(tables) == 1
        assert tables[0]["rows"] == 2
        assert tables[0]["columns"] == 3
        assert tables[0]["data"][0] == ["Name", "Age", "City"]
        assert tables[0]["data"][1] == ["Alice", "30", "NYC"]


# === XLSX Tests ===

class TestXlsxReader:
    def test_read_xlsx(self, sample_xlsx):
        text = read_xlsx(sample_xlsx)
        assert "Data" in text
        assert "Alice" in text
        assert "Summary" in text

    def test_get_xlsx_info(self, sample_xlsx):
        info = get_xlsx_info(sample_xlsx)
        assert info["sheets"] == 2
        assert info["sheet_details"][0]["name"] == "Data"
        assert info["sheet_details"][0]["rows"] == 3
        assert info["sheet_details"][0]["columns"] == 3
        assert info["sheet_details"][1]["name"] == "Summary"

    def test_get_xlsx_sheets_all(self, sample_xlsx):
        sheets = get_xlsx_sheets(sample_xlsx)
        assert len(sheets) == 2
        assert sheets[0]["sheet"] == "Data"
        assert sheets[0]["data"][0] == ["Name", "Score", "Grade"]
        assert sheets[0]["data"][1] == ["Alice", "95", "A"]

    def test_get_xlsx_sheets_filtered(self, sample_xlsx):
        sheets = get_xlsx_sheets(sample_xlsx, sheet_name="Summary")
        assert len(sheets) == 1
        assert sheets[0]["sheet"] == "Summary"
        assert sheets[0]["data"][0] == ["Total", "177"]

    def test_get_xlsx_sheets_max_rows(self, sample_xlsx):
        sheets = get_xlsx_sheets(sample_xlsx, max_rows=2)
        assert sheets[0]["rows_returned"] == 2
        assert sheets[0]["total_rows"] == 3


# === PPTX Tests ===

class TestPptxReader:
    def test_read_pptx(self, sample_pptx):
        text = read_pptx(sample_pptx)
        assert "Welcome" in text
        assert "Subtitle text" in text
        assert "Content Slide" in text

    def test_get_pptx_info(self, sample_pptx):
        info = get_pptx_info(sample_pptx)
        assert info["title"] == "Test Presentation"
        assert info["author"] == "Test Author"
        assert info["slides"] == 2
        assert info["total_shapes"] > 0

    def test_get_pptx_slides(self, sample_pptx):
        slides = get_pptx_slides(sample_pptx)
        assert len(slides) == 2
        assert slides[0]["index"] == 0
        assert "Welcome" in slides[0]["text_content"]
        assert slides[1]["index"] == 1
        assert "Content Slide" in slides[1]["text_content"]
